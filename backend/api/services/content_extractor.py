"""
content_extractor.py — Document text extraction for the semantic filesystem index.

Pulls searchable text out of common document types so fs_index can embed
real content instead of just filenames. Every extractor is defensive:
corrupt/encrypted/oversized files return None rather than raising, so a
single bad file never takes down an index rebuild or watcher event.

Logs: [CONTENT_EXTRACT_OK] [CONTENT_EXTRACT_SKIP] [CONTENT_EXTRACT_FAIL]
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# Files larger than this are never opened for content extraction — avoids
# multi-GB video/ISO/archive files masquerading under a supported extension.
MAX_EXTRACT_BYTES = 25 * 1024 * 1024  # 25 MB

# Max characters kept per file — enough for embedding + keyword search
# without bloating SQLite rows or the FAISS text cache.
MAX_TEXT_CHARS = 8000

TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".csv",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".json", ".yaml", ".yml",
    ".java", ".c", ".cpp", ".h", ".hpp", ".go", ".rs", ".sh", ".bash",
    ".html", ".css", ".sql", ".toml", ".ini", ".cfg", ".xml",
}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
PPTX_EXTS = {".pptx"}
XLSX_EXTS = {".xlsx", ".xlsm"}

SUPPORTED_EXTS = TEXT_EXTS | PDF_EXTS | DOCX_EXTS | PPTX_EXTS | XLSX_EXTS


def is_supported(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTS


def extract_text(path: Path) -> Optional[str]:
    """
    Best-effort text extraction. Returns None if unsupported, oversized,
    unreadable, or extraction fails for any reason.
    """
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        return None

    try:
        if path.stat().st_size > MAX_EXTRACT_BYTES:
            logger.debug("[CONTENT_EXTRACT_SKIP] path=%s reason=too_large", path)
            return None
    except OSError:
        return None

    try:
        if ext in TEXT_EXTS:
            text = _extract_plain_text(path)
        elif ext in PDF_EXTS:
            text = _extract_pdf(path)
        elif ext in DOCX_EXTS:
            text = _extract_docx(path)
        elif ext in PPTX_EXTS:
            text = _extract_pptx(path)
        elif ext in XLSX_EXTS:
            text = _extract_xlsx(path)
        else:
            return None
    except Exception as exc:  # noqa: BLE001 — a bad document must never break indexing
        logger.debug("[CONTENT_EXTRACT_FAIL] path=%s error=%s", path, exc)
        return None

    if not text:
        return None

    text = text.strip()[:MAX_TEXT_CHARS]
    if text:
        logger.debug("[CONTENT_EXTRACT_OK] path=%s chars=%d", path, len(text))
    return text or None


def _extract_plain_text(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read(MAX_TEXT_CHARS * 2)  # over-read slightly; caller truncates


def _extract_pdf(path: Path) -> str:
    import fitz  # PyMuPDF
    parts = []
    with fitz.open(str(path)) as doc:
        for page in doc:
            parts.append(page.get_text())
            if sum(len(p) for p in parts) >= MAX_TEXT_CHARS:
                break
    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    import docx
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            for row in ws.iter_rows(max_row=200, values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
                if sum(len(p) for p in parts) >= MAX_TEXT_CHARS:
                    return "\n".join(parts)
    finally:
        wb.close()
    return "\n".join(parts)
